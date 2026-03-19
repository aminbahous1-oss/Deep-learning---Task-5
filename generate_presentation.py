from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

PURPLE       = RGBColor(0x7C, 0x3A, 0xED)
LIGHT_PURPLE = RGBColor(0xED, 0xE9, 0xFE)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK         = RGBColor(0x1F, 0x15, 0x4B)
GRAY         = RGBColor(0x6B, 0x72, 0x80)
YELLOW       = RGBColor(0xFC, 0xD3, 0x4D)
BLUE         = RGBColor(0x60, 0xA5, 0xFA)
RED          = RGBColor(0xF8, 0x71, 0x71)
GREEN        = RGBColor(0x6E, 0xE7, 0xB7)

def blank_slide():
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=WHITE):
    bg_shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color
    bg_shape.line.fill.background()

def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=DARK,
                 align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_rect(slide, left, top, width, height, color=PURPLE, radius=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.4, PURPLE)
    add_text_box(slide, title, 0.4, 0.15, 12, 0.7, font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text_box(slide, subtitle, 0.4, 0.85, 12, 0.45, font_size=14, color=LIGHT_PURPLE, align=PP_ALIGN.LEFT)

def bullet_box(slide, items, left, top, width, height, font_size=16, color=DARK, bullet="•"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color

# ── Slide 1: Title ───────────────────────────────────────────────────────────
s = blank_slide()
bg(s, PURPLE)
add_rect(s, 0, 2.5, 13.33, 3.2, DARK)
add_text_box(s, "🧠 MindJournal", 0.5, 2.7, 12, 1.1, font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, "A Smart Mental Health Journal powered by Cognitive Services",
             0.5, 3.85, 12, 0.6, font_size=20, color=LIGHT_PURPLE, align=PP_ALIGN.CENTER)
add_text_box(s, "TIES4911 – Lecture 8  |  University of Jyväskylä",
             0.5, 6.7, 12, 0.5, font_size=13, color=LIGHT_PURPLE, align=PP_ALIGN.CENTER)

# ── Slide 2: Problem ─────────────────────────────────────────────────────────
s = blank_slide()
bg(s)
header_bar(s, "The Problem", "Why do we need a smart journal?")
problems = [
    "1 in 4 people struggle to recognize and track their own emotional patterns",
    "Traditional journaling gives no feedback or actionable insights",
    "Mental health check-ins are often skipped — they feel tedious",
    "People lack tools that bridge self-reflection and data-driven awareness",
]
bullet_box(s, problems, 0.5, 1.6, 12, 3.5, font_size=18)
add_rect(s, 0.5, 5.4, 12.3, 0.9, LIGHT_PURPLE)
add_text_box(s, "Goal: Make emotional self-awareness effortless and data-driven",
             0.7, 5.5, 12, 0.7, font_size=17, bold=True, color=PURPLE)

# ── Slide 3: Solution ────────────────────────────────────────────────────────
s = blank_slide()
bg(s)
header_bar(s, "What is MindJournal?", "Our solution")
add_text_box(s, "A web-based smart journal that:", 0.5, 1.6, 12, 0.5, font_size=18, bold=True, color=DARK)
features = [
    "Lets you speak or type how you feel — no friction",
    "Analyzes your emotions from both your text and your face (webcam)",
    "Tracks your mood over time with interactive visual charts",
    "Requires zero mental health expertise to use",
    "Works instantly — no sign-up, runs in your browser",
]
bullet_box(s, features, 0.5, 2.2, 12, 3.5, font_size=17)

# ── Slide 4: Cognitive Services ──────────────────────────────────────────────
s = blank_slide()
bg(s)
header_bar(s, "Cognitive Services Used", "Multi-provider integration")

cards = [
    ("IBM Watson NLU", "Emotion + Sentiment\nfrom text", PURPLE),
    ("Web Speech API\n(Browser)", "Voice-to-text\ninput", BLUE),
    ("face-api.js\n(Open Source ML)", "Real-time webcam\nemotion detection", GREEN),
    ("Chart.js", "Mood trend\nvisualization", YELLOW),
]
for i, (title, desc, color) in enumerate(cards):
    x = 0.4 + i * 3.15
    add_rect(s, x, 1.6, 2.9, 2.5, color)
    add_text_box(s, title, x+0.1, 1.7, 2.7, 0.9, font_size=15, bold=True, color=WHITE)
    add_text_box(s, desc, x+0.1, 2.65, 2.7, 1.2, font_size=14, color=WHITE)

add_text_box(s, "Providers: IBM  ·  Browser Native API  ·  Open Source",
             0.5, 4.3, 12, 0.5, font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ── Slide 5: Architecture ────────────────────────────────────────────────────
s = blank_slide()
bg(s)
header_bar(s, "System Architecture", "How it all connects")

boxes = [
    (0.4,  2.8, "User\nSpeaks / Types", PURPLE),
    (3.2,  2.0, "Speech API\n(Browser)", BLUE),
    (3.2,  3.8, "Webcam\n(face-api.js)", GREEN),
    (6.2,  2.0, "Flask\nBackend", DARK),
    (6.2,  3.8, "IBM Watson\nNLU", PURPLE),
    (9.5,  2.0, "localStorage\n+ Charts", YELLOW),
]
for (x, y, label, color) in boxes:
    add_rect(s, x, y, 2.4, 1.1, color)
    tcol = DARK if color == YELLOW else WHITE
    add_text_box(s, label, x+0.05, y+0.15, 2.3, 0.8, font_size=14, bold=True, color=tcol, align=PP_ALIGN.CENTER)

arrows = [
    "User → Speech API & Webcam  →  Flask Backend & IBM Watson  →  localStorage  →  Dashboard Charts"
]
add_text_box(s, arrows[0], 0.4, 5.3, 12.5, 0.6, font_size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ── Slide 6: Journal Page ────────────────────────────────────────────────────
s = blank_slide()
bg(s)
header_bar(s, "Feature 1 — Journal Entry Page", "app: http://localhost:5000")

add_rect(s, 0.4, 1.6, 5.8, 5.3, LIGHT_PURPLE)
add_text_box(s, "Journal Entry", 0.6, 1.7, 5.4, 0.4, font_size=14, bold=True, color=PURPLE)
add_rect(s, 0.5, 2.15, 5.6, 2.0, RGBColor(0xFF,0xFF,0xFF))
add_text_box(s, "How are you feeling today?...", 0.6, 2.25, 5.4, 1.8, font_size=13, color=GRAY)
add_rect(s, 0.5, 4.3, 2.7, 0.5, PURPLE)
add_text_box(s, "🎤 Start Recording", 0.55, 4.35, 2.6, 0.4, font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s, 3.3, 4.3, 2.7, 0.5, RGBColor(0x10,0xB9,0x81))
add_text_box(s, "📸 Capture Mood", 3.35, 4.35, 2.6, 0.4, font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s, 0.5, 5.0, 5.6, 0.55, PURPLE)
add_text_box(s, "Analyze Entry", 0.55, 5.05, 5.5, 0.45, font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s, 0.5, 5.7, 5.6, 0.8, RGBColor(0xFF,0xF9,0xC4))
add_text_box(s, "⚠️  Using demo data — add Watson API key for real analysis",
             0.6, 5.75, 5.4, 0.7, font_size=11, color=RGBColor(0x92,0x40,0x0E))

add_rect(s, 6.8, 1.6, 5.9, 5.3, LIGHT_PURPLE)
add_text_box(s, "Webcam Feed", 7.0, 1.7, 5.5, 0.4, font_size=14, bold=True, color=PURPLE)
add_rect(s, 6.9, 2.15, 5.7, 2.3, DARK)
add_text_box(s, "[ live webcam ]", 6.95, 3.0, 5.6, 0.6, font_size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text_box(s, "Face Emotion: 😊 Joy  (82%)", 7.0, 4.55, 5.5, 0.45, font_size=14, bold=True, color=PURPLE)

emotions = [("Joy", 0.75, YELLOW), ("Sadness", 0.10, BLUE), ("Anger", 0.05, RED), ("Fear", 0.06, RGBColor(0xA7,0x8B,0xFA))]
for i, (name, val, col) in enumerate(emotions):
    y = 5.1 + i * 0.35
    add_text_box(s, name, 7.0, y, 1.2, 0.32, font_size=11, color=DARK)
    add_rect(s, 8.3, y+0.05, val * 3.5, 0.22, col)
    add_text_box(s, f"{int(val*100)}%", 8.35 + val*3.5, y, 0.6, 0.3, font_size=10, color=GRAY)

# ── Slide 7: Dashboard ───────────────────────────────────────────────────────
s = blank_slide()
bg(s)
header_bar(s, "Feature 2 — Mood Dashboard", "app: http://localhost:5000/dashboard")

stat_labels = ["Total Entries", "Avg Joy", "Avg Sadness", "Top Mood"]
stat_vals   = ["12", "68%", "14%", "😊 Joy"]
stat_colors = [PURPLE, YELLOW, BLUE, GREEN]
for i, (lbl, val, col) in enumerate(zip(stat_labels, stat_vals, stat_colors)):
    x = 0.4 + i * 3.15
    add_rect(s, x, 1.6, 2.9, 1.1, col)
    tcol = DARK if col == YELLOW else WHITE
    add_text_box(s, val,  x+0.1, 1.65, 2.7, 0.55, font_size=24, bold=True, color=tcol, align=PP_ALIGN.CENTER)
    add_text_box(s, lbl,  x+0.1, 2.2,  2.7, 0.4,  font_size=12, color=tcol, align=PP_ALIGN.CENTER)

add_rect(s, 0.4, 2.9, 6.0, 3.5, LIGHT_PURPLE)
add_text_box(s, "📈 Emotion Trends (last 14 entries)", 0.6, 2.95, 5.6, 0.4, font_size=13, bold=True, color=PURPLE)
add_text_box(s, "[ Line Chart: Joy / Sadness / Anger / Fear / Disgust over time ]",
             0.6, 3.5, 5.6, 2.5, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

add_rect(s, 6.9, 2.9, 5.9, 3.5, LIGHT_PURPLE)
add_text_box(s, "📊 Average Emotion Distribution", 7.1, 2.95, 5.5, 0.4, font_size=13, bold=True, color=PURPLE)
add_text_box(s, "[ Bar Chart: average values across all entries ]",
             7.1, 3.5, 5.5, 2.5, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ── Slide 8: Future Work ─────────────────────────────────────────────────────
s = blank_slide()
bg(s)
header_bar(s, "Possible Extensions", "Where MindJournal could go next")
future = [
    "AWS Comprehend — deeper sentiment analysis and entity detection",
    "Azure Translator — multi-language journal support",
    "User accounts with cloud storage (Firebase / Supabase)",
    "Daily reminder notifications via email or push",
    "Therapist sharing mode — export mood reports as PDF",
    "Mobile app version (React Native)",
    "Integration with wearables (heart rate + mood correlation)",
]
bullet_box(s, future, 0.5, 1.6, 12, 5.0, font_size=17)

# ── Slide 9: Conclusion ──────────────────────────────────────────────────────
s = blank_slide()
bg(s, PURPLE)
add_text_box(s, "Summary", 0.5, 0.8, 12, 0.8, font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
points = [
    "✅  Built a fully working Smart Mental Health Journal prototype",
    "✅  Combined 3 cognitive service sources (IBM Watson, Browser API, face-api.js)",
    "✅  Works immediately with demo mode — real analysis enabled by one API key",
    "✅  Demonstrates real-world value in the Healthcare / Wellness domain",
    "✅  Extensible to AWS, Azure, and mobile platforms",
]
bullet_box(s, points, 1.0, 1.9, 11, 4.0, font_size=18, color=WHITE, bullet="")
add_text_box(s, "TIES4911 – University of Jyväskylä", 0.5, 6.7, 12, 0.5,
             font_size=13, color=LIGHT_PURPLE, align=PP_ALIGN.CENTER)

prs.save("MindJournal_Presentation.pptx")
print("Saved: MindJournal_Presentation.pptx")
